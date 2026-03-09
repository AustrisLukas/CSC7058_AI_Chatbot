# from helpers import helpers
from pathlib import Path
import re
from pypdf import PdfReader
import docx
import pandas as pd
import logging
from .exceptions import DocumentExtractionError
from pptx import Presentation
import json
import requests
import trafilatura
from trafilatura.metadata import extract_metadata


logger = logging.getLogger(__name__)


def detect_file_type(file):
    EXTENSION_MAP = {
        ".pdf": "pdf",
        ".txt": "text",
        ".docx": "word",
        ".doc": "word",
        ".pptx": "powerpoint",
        ".ppt": "powerpoint",
        ".xlsx": "excel",
        ".xls": "excel",
        ".csv": "csv",
    }
    MIME_MAP = {
        "application/pdf": "pdf",
        "text/plain": "text",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "word",
        "application/msword": "word",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "powerpoint",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "excel",
        "text/csv": "csv",
    }

    extension = Path(file.name).suffix.lower()
    mime = file.type
    extension_type = EXTENSION_MAP.get(extension)
    mime_type = MIME_MAP.get(mime)

    if extension_type and mime_type and extension_type == mime_type:
        return extension_type
    if extension_type and not mime_type:
        return extension_type

    raise ValueError("Unsupported or mismatch file type")


def extract_pdf(file):
    try:
        reader = PdfReader(file)
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        # .pdf SPECIFIC CLEAN UP
        text = text.replace("\n", " ")  # MERGE LINES
        text = re.sub(r"\s", " ", text)  # COLLAPSE MULTIPLE SPACES
        text = re.sub(r"Page \d+ of \d+", "", text)  # REMOVE PAGE NUMBERS
        text = text.strip()

        if not text:
            print(text)
            raise DocumentExtractionError("No text could be extracted from the PDF")
        return text
    except Exception as e:
        raise DocumentExtractionError(f"Document Processing Error: {e}") from e


def extract_docx(doc):
    try:
        doc = docx.Document(doc)
        text = ""
        for p in doc.paragraphs:
            if p.text:
                text += p.text + " "

        text = text.replace("\n", " ")  # MERGE LINES
        text = re.sub(r"\s", " ", text)  # COLLAPSE MULTIPLE SPACES

        if not text:
            raise DocumentExtractionError(
                "No text could be extracted from the document"
            )
        return text
    except Exception as e:
        raise DocumentExtractionError(f"Document Processing Error: {e}") from e


def extract_csv(file):

    try:
        text = ""
        file.seek(0)  # set reader pointer to 0
        df = pd.read_csv(file, header=None)
        df = df.fillna("")  # fill empty cells with string
        df = df.astype(str)  # ensure all data types are string
        text = " ".join(df.agg(" ".join, axis=1))  # combine in one string

        if not text:
            raise DocumentExtractionError(
                "No text could be extracted from the CSV file."
            )
        return text
    except Exception as e:
        raise DocumentExtractionError(f"Document Processing Error: {e}") from e


def extract_excel(file):
    try:
        text = ""
        file.seek(0)  # reset file read to start
        df = pd.read_excel(file, header=None)
        df = df.fillna("")  # fill empty cells with string
        df = df.astype(str)  # ensure all data types are string
        text = " ".join(df.agg(" ".join, axis=1))  # combine in one string
        if not text:
            raise DocumentExtractionError(
                "No text could be extracted from the excel file"
            )
        return text
    except Exception as e:
        raise DocumentExtractionError(f"Document Processing Error: {e}") from e


def extract_pptx(file):

    prs = Presentation(file)

    text = ""

    try:
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text + " "
        if not text:
            raise DocumentExtractionError(
                "No text could be extracted from the power point file"
            )
        return text
    except Exception as e:
        raise DocumentExtractionError(f"Document Processing Error: {e}") from e


def extract_from_url(url: str):
    try:
        response = requests.get(url)
        metadata = extract_metadata(response.text)
        text = trafilatura.extract(
            response.text,
            include_links=True,
            include_images=True,
            include_tables=True,
            favor_precision=False,
            favor_recall=True,
            include_comments=False,
            deduplicate=False,
            output_format="markdown",
        )

        extr_document = f"""
        Title: {metadata.title}
        Website: {metadata.sitename}
        URL: {metadata.url}

        {clean_extracted_text(text)}
        """

        return extr_document
    except Exception as e:
        logger.error(f"Error while extracting text from provided URL: {e}")
        raise DocumentExtractionError(
            f"An error has occurred while extracting text from URL: {e}"
        ) from e


def clean_extracted_text(text):

    if not text:
        return ""

    # Normalise line endings
    text = text.replace("\r\n", "\n")
    # Remove separation characters
    text = text.replace("|", " ")
    # Remove bullet point characters
    text = re.sub(r"[•▪◦●■*]", " ", text)
    # Normalise multiple spaces
    text = re.sub(r"[ \t]+", " ", text)
    # Normalise excessive newlines (max two)
    text = re.sub(r"\n\s*\n\s*\n+", "\n\n", text)
    # Remove leading/trailing whitespace
    text = text.strip()

    return text


def parse_model_json(text: str) -> dict:
    s = text.strip()

    # remove ```json ... ``` or ``` ... ```
    if s.startswith("```"):
        lines = s.splitlines()
        if lines and lines[0].startswith("```"):
            lines = lines[1:]
        if lines and lines[-1].strip() == "```":
            lines = lines[:-1]
        s = "\n".join(lines).strip()

    return json.loads(s)


def load_json(raw_json):

    try:
        data = json.loads(raw_json)
        return data
    except json.JSONDecodeError:
        data = {
            "answer": "Model returned invalid JSON.",
            "self_score": 0,
            "reason": "Parsing failure.",
            "references": [],
        }
        return data
